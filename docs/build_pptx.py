"""
Proofolio 발표 슬라이드 빌더 (.pptx).

실행:
    python3 docs/build_pptx.py

산출:
    docs/Proofolio_발표.pptx  (Keynote / PowerPoint / Google Slides 호환)

톤: 쿨톤 라이트 + 에메랄드(verified) 액센트 (앱 디자인과 동일).
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ─── 디자인 토큰 ────────────────────────────────────────────────
BG = RGBColor(0xF7, 0xF8, 0xFA)            # paper-50
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
HAIRLINE = RGBColor(0xE6, 0xE8, 0xEE)      # ink-100
INK_950 = RGBColor(0x0B, 0x12, 0x20)       # 본문/제목
INK_700 = RGBColor(0x3F, 0x47, 0x57)       # 보조
INK_500 = RGBColor(0x6B, 0x72, 0x80)       # 뮤트
ACCENT = RGBColor(0x05, 0x96, 0x69)        # trust-600 (emerald)
ACCENT_LIGHT = RGBColor(0x10, 0xB9, 0x81)  # trust-500

FONT = "Pretendard"   # 없으면 Apple SD Gothic Neo로 자동 대체
FONT_MONO = "Menlo"

# ─── 레이아웃 헬퍼 ──────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def text(slide, left, top, width, height, content,
         *, size=18, bold=False, color=INK_950, font=FONT,
         align=PP_ALIGN.LEFT, line_spacing=1.3):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    lines = content.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = color
    return box


def eyebrow(slide, top, content, color=ACCENT):
    return text(
        slide, Inches(0.8), top, Inches(11), Inches(0.35),
        content.upper(), size=11, bold=True, color=color,
    )


def hairline(slide, top, color=HAIRLINE):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), top, Inches(11.7), Emu(8000),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def placeholder(slide, left, top, width, height, label, sublabel=""):
    """미디어(영상·이미지) 자리. Keynote에서 .mov/.png 드래그 시 이 박스가 교체됨."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF1, 0xF3, 0xF7)
    shape.line.color.rgb = HAIRLINE
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(20)
    tf.margin_top = tf.margin_bottom = Pt(20)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = INK_500
    run.font.name = FONT

    if sublabel:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(6)
        r2 = p2.add_run()
        r2.text = sublabel
        r2.font.size = Pt(10)
        r2.font.color.rgb = INK_500
        r2.font.name = FONT

    return shape


def pill(slide, left, top, content, *, fg=ACCENT, bg_alpha=0xE9):
    # 작은 둥근 알약 (참고/링크용)
    width = Inches(0.05 + 0.085 * len(content))
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.36)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xE7, 0xF6, 0xEE)
    rect.line.color.rgb = RGBColor(0xC8, 0xEA, 0xD9)
    rect.line.width = Pt(0.75)
    tf = rect.text_frame
    tf.margin_left = tf.margin_right = Pt(8)
    tf.margin_top = tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = content
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.name = FONT
    run.font.color.rgb = fg
    return rect


# ─── Slide 1: 표지 ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s)
text(s, Inches(0.8), Inches(0.8), Inches(4), Inches(0.4),
     "PROOFOLIO", size=11, bold=True, color=ACCENT)

text(s, Inches(0.8), Inches(2.4), Inches(12), Inches(1.5),
     "검증 가능한\n디지털 증명서.",
     size=64, bold=True, color=INK_950, line_spacing=1.05)

text(s, Inches(0.8), Inches(4.7), Inches(12), Inches(1),
     "PDF가 아니라 발급기관의 온체인 서명을 믿는다.",
     size=20, color=INK_700)

# 하단
text(s, Inches(0.8), Inches(6.45), Inches(11), Inches(0.4),
     "이름 · 학번 · 2026.06", size=12, color=INK_500)
text(s, Inches(0.8), Inches(6.85), Inches(11), Inches(0.4),
     "https://proofolio.pages.dev", size=12, color=ACCENT)


# ─── Slide 2: Problem ─────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s)
eyebrow(s, Inches(0.6), "01 · Problem")

text(s, Inches(0.8), Inches(1.2), Inches(12), Inches(1.5),
     "가짜 이력서,\n위조 수료증.",
     size=52, bold=True, color=INK_950, line_spacing=1.1)

hairline(s, Inches(4.4))

text(s, Inches(0.8), Inches(4.8), Inches(12), Inches(0.6),
     "· 채용 담당자는 제출된 PDF를 그대로 믿거나",
     size=20, color=INK_700)
text(s, Inches(0.8), Inches(5.4), Inches(12), Inches(0.6),
     "· 발급기관에 일일이 전화해서 확인해야 한다.",
     size=20, color=INK_700)

text(s, Inches(0.8), Inches(6.5), Inches(12), Inches(0.6),
     "매년 반복되는, 돈이 걸린 진짜 문제.",
     size=20, bold=True, color=INK_950)


# ─── Slide 3: Solution & Why blockchain ───────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s)
eyebrow(s, Inches(0.6), "02 · Solution · Why blockchain")

text(s, Inches(0.8), Inches(1.2), Inches(12), Inches(1.8),
     "PDF가 아니라\n발급기관의 온체인 서명을 믿는다.",
     size=40, bold=True, color=INK_950, line_spacing=1.15)

hairline(s, Inches(4.2))

# 좌: 기존 DB
text(s, Inches(0.8), Inches(4.6), Inches(5.5), Inches(0.35),
     "기존 DB", size=12, bold=True, color=INK_500)
text(s, Inches(0.8), Inches(5.0), Inches(5.5), Inches(2),
     "운영자가 기록을 몰래 수정 가능 →\n검증이 결국 '운영자 신뢰 게임'으로 회귀.",
     size=15, color=INK_700, line_spacing=1.5)

# 우: 블록체인
text(s, Inches(6.8), Inches(4.6), Inches(5.7), Inches(0.35),
     "블록체인", size=12, bold=True, color=ACCENT)
text(s, Inches(6.8), Inches(5.0), Inches(5.7), Inches(0.45),
     "✓  운영자도 못 바꾼다.",
     size=15, color=INK_950)
text(s, Inches(6.8), Inches(5.5), Inches(5.7), Inches(0.45),
     "✓  Etherscan에서 누구나 독립 검증.",
     size=15, color=INK_950)
text(s, Inches(6.8), Inches(6.0), Inches(5.7), Inches(0.45),
     "✓  파일 한 글자만 바뀌어도 해시 불일치.",
     size=15, color=INK_950)


# ─── Slide 4: Demo (영상 + 라이브 전환) ───────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s)
eyebrow(s, Inches(0.6), "03 · Demo")

text(s, Inches(0.8), Inches(1.15), Inches(12), Inches(0.7),
     "발급 영상 → 라이브 검증",
     size=32, bold=True, color=INK_950, line_spacing=1.1)

# 영상 placeholder — Keynote 열어서 .mov 파일을 이 박스로 드래그
VIDEO_W = Inches(9)
VIDEO_H = Inches(4.5)
VIDEO_LEFT = Inches((13.333 - 9) / 2)  # 가운데 정렬
VIDEO_TOP = Inches(2.2)
placeholder(
    s, VIDEO_LEFT, VIDEO_TOP, VIDEO_W, VIDEO_H,
    "▶  발급 흐름 영상",
    "이 자리에 .mov 파일을 드래그 (자동재생, 음소거, 15~20초)",
)

text(s, Inches(0.8), Inches(6.95), Inches(12), Inches(0.4),
     "영상 자동재생 후 → 브라우저로 전환해 라이브 검증",
     size=12, color=INK_500, align=PP_ALIGN.CENTER)


# ─── Slide 5: 그 외 화면 (스샷 3장) ───────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s)
eyebrow(s, Inches(0.6), "04 · 그 외 화면")

text(s, Inches(0.8), Inches(1.15), Inches(12), Inches(0.7),
     "관리자 · 발급기관 · 보유자",
     size=32, bold=True, color=INK_950, line_spacing=1.1)

text(s, Inches(0.8), Inches(1.95), Inches(12), Inches(0.5),
     "검증 외에도 다음 화면들이 같은 디자인 시스템으로 구현되어 있습니다.",
     size=14, color=INK_700, line_spacing=1.4)

# 스샷 placeholder 3개 — 가로 배치 (모든 위치 계산은 inch 단위로)
SHOT_W_IN = 3.85
SHOT_H_IN = 3.0
SHOT_GAP_IN = 0.25
SHOT_TOP_IN = 3.0
shot_total_w_in = SHOT_W_IN * 3 + SHOT_GAP_IN * 2
shot_left_start_in = (13.333 - shot_total_w_in) / 2

screens = [
    ("관리자 · /admin", "발급기관 등록 · 활성 관리"),
    ("발급 · /issue", "보유자에게 증명서 발급"),
    ("보유자 · /credentials", "받은 증명서 보기"),
]
for i, (label, sub) in enumerate(screens):
    x_in = shot_left_start_in + (SHOT_W_IN + SHOT_GAP_IN) * i
    placeholder(s, Inches(x_in), Inches(SHOT_TOP_IN),
                Inches(SHOT_W_IN), Inches(SHOT_H_IN),
                label, "스샷을 이 자리에 드래그")
    text(s, Inches(x_in), Inches(SHOT_TOP_IN + SHOT_H_IN + 0.15),
         Inches(SHOT_W_IN), Inches(0.4),
         sub, size=11, color=INK_500, align=PP_ALIGN.CENTER, line_spacing=1.3)


# ─── Slide 6: Architecture & Engineering ──────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s)
eyebrow(s, Inches(0.6), "05 · Architecture & Engineering")

text(s, Inches(0.8), Inches(1.2), Inches(12), Inches(1),
     "단순한 데모가 아닙니다.",
     size=36, bold=True, color=INK_950, line_spacing=1.1)

# 다이어그램 — 3단계 흐름, 평이한 한국어만
arch_lines = [
    "발급기관    →    블록체인에 영구 기록    →    보유자 지갑",
    "",
    "                              ↑",
    "             검증자가 지갑 없이도 누구나 진위 확인",
]
text(s, Inches(0.8), Inches(2.7), Inches(12), Inches(2.4),
     "\n".join(arch_lines),
     size=17, color=INK_700, line_spacing=1.6, align=PP_ALIGN.CENTER)

hairline(s, Inches(5.35))

points = [
    "증명서는 양도가 불가능한 형태로 발급 — 사고팔 수 없음",
    "자동 테스트 17개로 정상 · 비정상 · 공격 시나리오까지 검증",
    "코드를 Etherscan에 공개 — 저희를 안 믿어도 코드는 검증 가능",
    "한 가지 한계 — '기관이 발급했다'는 보장하지만 '기관이 좋은가'는 별도 검토",
]
y = Inches(5.6)
for p in points:
    text(s, Inches(0.8), y, Inches(12), Inches(0.4),
         "·  " + p, size=14, color=INK_950)
    y += Inches(0.4)


# ─── Slide 6: Vibe coding & 마무리 ────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s)
eyebrow(s, Inches(0.6), "06 · Vibe Coding · 노하우")

text(s, Inches(0.8), Inches(1.2), Inches(12), Inches(1),
     "AI로 어떻게 만들었나.",
     size=36, bold=True, color=INK_950, line_spacing=1.1)

vibe = [
    ("01  문서 먼저, 코드 나중",
     "PRD / SRS 1000줄을 AI와 같이 작성. AGENTS.md 로 규칙 박제."),
    ("02  Phase로 잘게 (10단계)",
     "단계마다 검증·커밋·1줄 회고. 커밋 히스토리가 곧 워크플로 증거."),
    ("03  AI 코드 ≠ 신뢰",
     "Hardhat 테스트 통과해야 신뢰. 실제로 AI deploy 스크립트 버그 잡아냄."),
    ("04  토큰 분업",
     "기획·디자인·보안 리뷰 = Opus / 반복 구현 = Codex."),
]
y = Inches(2.7)
for title, body in vibe:
    text(s, Inches(0.8), y, Inches(12), Inches(0.4),
         title, size=15, bold=True, color=ACCENT)
    text(s, Inches(0.8), y + Inches(0.4), Inches(12), Inches(0.5),
         body, size=14, color=INK_700, line_spacing=1.35)
    y += Inches(1.05)

# 마무리 한 줄
text(s, Inches(0.8), Inches(6.85), Inches(12), Inches(0.5),
     "AI는 빠른 손, 사람은 기준과 검증.   감사합니다.",
     size=14, color=INK_500)


# ─── 저장 ─────────────────────────────────────────────────────
out = "docs/Proofolio_발표.pptx"
prs.save(out)
print(f"✓ Saved: {out}  ({len(prs.slides)} slides, 16:9)")
