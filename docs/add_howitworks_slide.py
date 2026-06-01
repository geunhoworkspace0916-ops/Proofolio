"""
사용자 PPT(~/Downloads/Proofolio_발표.pptx)의 슬라이드 3(Solution) 뒤에
"작동 원리" 슬라이드 1장을 삽입한다. — 해시 + 지갑 서명(누가 발급했나) 동시 설명.

실행: python3 docs/add_howitworks_slide.py
결과: ~/Downloads/Proofolio_발표.pptx 가 갱신됨 (원본은 .bak 으로 백업)
"""

import copy
import shutil
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

PATH = "/Users/geunho/Downloads/Proofolio_발표.pptx"

# ── 사용자 PPT와 동일한 디자인 토큰 ──
BG = RGBColor(0xF7, 0xF8, 0xFA)
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
HAIRLINE = RGBColor(0xE6, 0xE8, 0xEE)
INK_950 = RGBColor(0x0B, 0x12, 0x20)
INK_700 = RGBColor(0x3F, 0x47, 0x57)
INK_500 = RGBColor(0x6B, 0x72, 0x80)
ACCENT = RGBColor(0x05, 0x96, 0x69)
FONT = "Pretendard"
FONT_MONO = "Menlo"

shutil.copyfile(PATH, PATH + ".bak")
prs = Presentation(PATH)
BLANK = prs.slide_layouts[6]


def text(slide, left, top, width, height, content, *, size=18, bold=False,
         color=INK_950, font=FONT, align=PP_ALIGN.LEFT, line_spacing=1.3):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(content.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = font
        r.font.color.rgb = color
    return box


def card(slide, left, top, width, height, *, fill=SURFACE, border=HAIRLINE):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border
    sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


# ── 새 슬라이드 생성 ──
s = prs.slides.add_slide(BLANK)
s.background.fill.solid()
s.background.fill.fore_color.rgb = BG

# eyebrow
text(s, Inches(0.8), Inches(0.6), Inches(11), Inches(0.35),
     "03 · HOW IT WORKS · 작동 원리", size=11, bold=True, color=ACCENT)

# 제목
text(s, Inches(0.8), Inches(1.1), Inches(12), Inches(0.9),
     "위조를 어떻게 잡아내는가.", size=36, bold=True, color=INK_950)

# ── 좌측 카드: 해시 (무엇이 진짜 파일인가) ──
LX, LY, LW, LH = Inches(0.8), Inches(2.25), Inches(5.75), Inches(4.4)
card(s, LX, LY, LW, LH)
text(s, LX + Inches(0.35), LY + Inches(0.3), LW - Inches(0.7), Inches(0.4),
     "① 해시 — 이 파일이 진짜 원본인가", size=16, bold=True, color=INK_950)

text(s, LX + Inches(0.35), LY + Inches(0.95), LW - Inches(0.7), Inches(0.9),
     "파일을 '해시'에 넣으면\n그 파일만의 고유 번호가 나온다.",
     size=14, color=INK_700, line_spacing=1.45)

# 해시 도식 (모노)
text(s, LX + Inches(0.35), LY + Inches(2.0), LW - Inches(0.7), Inches(0.5),
     "원본파일 ─▶ 0x7a3f…e2", size=13, color=ACCENT, font=FONT_MONO)

text(s, LX + Inches(0.35), LY + Inches(2.6), LW - Inches(0.7), Inches(1.5),
     "· 같은 파일 → 항상 같은 번호\n"
     "· 한 글자만 달라도 → 완전 다른 번호\n"
     "· 발급할 때 그 번호를 블록체인에 저장",
     size=13, color=INK_700, line_spacing=1.6)

# ── 우측 카드: 지갑 서명 (누가 발급했나) ──
RX, RY, RW, RH = Inches(6.85), Inches(2.25), Inches(5.65), Inches(4.4)
card(s, RX, RY, RW, RH)
text(s, RX + Inches(0.35), RY + Inches(0.3), RW - Inches(0.7), Inches(0.4),
     "② 지갑 서명 — 누가 발급했나", size=16, bold=True, color=INK_950)

text(s, RX + Inches(0.35), RY + Inches(0.95), RW - Inches(0.7), Inches(0.9),
     "발급기관마다 고유한 지갑 주소가\n신원이 된다. 위조 불가능.",
     size=14, color=INK_700, line_spacing=1.45)

text(s, RX + Inches(0.35), RY + Inches(2.0), RW - Inches(0.7), Inches(0.5),
     "발급기관 = 0xABCD…1234", size=13, color=ACCENT, font=FONT_MONO)

text(s, RX + Inches(0.35), RY + Inches(2.6), RW - Inches(0.7), Inches(1.5),
     "· 그 지갑만 자기 이름으로 발급 가능\n"
     "· 아이디·비밀번호 없이 주소가 곧 도장\n"
     "· 관리자가 등록한 기관만 발급 권한",
     size=13, color=INK_700, line_spacing=1.6)

# ── 하단 결론 한 줄 ──
text(s, Inches(0.8), Inches(6.85), Inches(12), Inches(0.5),
     "→ '등록된 기관(②)이 + 이 원본 파일(①)을' 발급했다 = 위조 불가능한 증명.",
     size=15, bold=True, color=INK_950)


# ── 새 슬라이드를 3번(index 2) 뒤 = 4번째 위치로 이동 ──
xml_slides = prs.slides._sldIdLst
slides = list(xml_slides)
new = slides[-1]              # 방금 추가된 마지막
xml_slides.remove(new)
xml_slides.insert(3, new)     # 0-based: index 3 = 4번째 슬라이드

prs.save(PATH)
print(f"✓ '작동 원리' 슬라이드를 4번째에 삽입 완료 ({len(prs.slides)}장)")
print(f"  원본 백업: {PATH}.bak")
